from __future__ import annotations

import io
import os
import re
import time
from typing import Dict, List, Type

SUPPORTED_EXTENSIONS: List[str] = [
    ".pdf",
    ".docx",
    ".txt",
    ".pptx",
    ".md",
    ".png",
    ".jpg",
    ".jpeg",
    ".tiff",
    ".tif",
    ".bmp",
]


class BaseExtractor:
    """Base extractor class. Subclasses implement `extract(path) -> str`."""

    def extract(self, path: str) -> str:
        raise NotImplementedError()


_registry: Dict[str, Type[BaseExtractor]] = {}


def register_extractor(exts: List[str]):
    def _decorator(cls: Type[BaseExtractor]):
        for e in exts:
            _registry[e.lower()] = cls
        return cls

    return _decorator


def extract_text_from_file(path: str) -> str:
    """Dispatch to the appropriate extractor based on file extension.

    Raises ValueError if extension is not supported.
    """
    _, ext = os.path.splitext(path)
    ext = ext.lower()
    if ext not in _registry:
        raise ValueError(f"Unsupported file extension: {ext}")
    extractor_cls = _registry[ext]
    extractor = extractor_cls()
    return extractor.extract(path)


@register_extractor([".txt", ".md"])
class TextExtractor(BaseExtractor):
    def extract(self, path: str) -> str:
        # Read bytes and try to decode safely
        with open(path, "rb") as f:
            data = f.read()
        # Try utf-8 then fallback to latin-1
        try:
            return data.decode("utf-8")
        except Exception:
            return data.decode("latin-1", errors="ignore")


@register_extractor([".docx"])
class DocxExtractor(BaseExtractor):
    def extract(self, path: str) -> str:
        try:
            from docx import Document
        except Exception as e:
            raise RuntimeError("Missing dependency: python-docx (install via requirements)") from e

        doc = Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text]
        return "\n".join(paragraphs)


@register_extractor([".pptx"])
class PptxExtractor(BaseExtractor):
    def extract(self, path: str) -> str:
        try:
            from pptx import Presentation
        except Exception as e:
            raise RuntimeError("Missing dependency: python-pptx (install via requirements)") from e

        prs = Presentation(path)
        texts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text
                    if txt:
                        texts.append(txt)
        return "\n".join(texts)


def _text_quality_score(text: str) -> float:
    if not text:
        return 0.0

    words = re.findall(r"\b\w+\b", text)
    if not words:
        return 0.0

    alpha_ratio = sum(c.isalpha() for c in text) / max(len(text), 1)
    long_words = [w for w in words if len(w) >= 4]
    bad_words = sum(1 for word in long_words if sum(c.isalpha() for c in word) / len(word) < 0.7)
    bad_ratio = bad_words / max(len(long_words), 1)
    return alpha_ratio * (1.0 - bad_ratio)


def _is_low_quality_text(text: str) -> bool:
    return _text_quality_score(text) < 0.65


def _preprocess_image(img):
    try:
        from PIL import Image, ImageOps
    except Exception:
        return img

    img = img.convert("L")
    img = ImageOps.autocontrast(img)
    if img.width < 1200:
        img = img.resize((img.width * 2, img.height * 2), Image.LANCZOS)
    return img


def _ocr_image(img):
    img = _preprocess_image(img)
    candidates: List[str] = []

    try:
        import pytesseract
    except Exception:
        pytesseract = None

    if pytesseract is not None:
        try:
            text = pytesseract.image_to_string(
                img,
                lang="eng",
                config="--oem 3 --psm 6",
            )
            if text and text.strip():
                candidates.append(text)
        except Exception:
            pass

    try:
        import easyocr
        import numpy as np
    except Exception as e:
        if candidates:
            return max(candidates, key=_text_quality_score)
        raise RuntimeError(
            "OCR requires Tesseract or easyocr. Install the Tesseract binary or install easyocr."
        ) from e

    try:
        reader = easyocr.Reader(["en"], gpu=False, verbose=False)
        results = reader.readtext(np.array(img))
        easy_text = "\n".join([text for _, text, _ in results])
        if easy_text and easy_text.strip():
            candidates.append(easy_text)
    except Exception:
        pass

    if not candidates:
        return ""

    return max(candidates, key=_text_quality_score)


@register_extractor([".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp"])
class ImageExtractor(BaseExtractor):
    def extract(self, path: str) -> str:
        try:
            from PIL import Image
        except Exception as e:
            raise RuntimeError(
                "Missing dependency: Pillow is required for image OCR"
            ) from e

        with Image.open(path) as img:
            return _ocr_image(img)


@register_extractor([".pdf"])
class PdfExtractor(BaseExtractor):
    def _ocr_page(self, page) -> str:
        try:
            import pytesseract
            from PIL import Image
        except Exception as e:
            raise RuntimeError(
                "Missing dependency: pytesseract and Pillow are required for OCR support"
            ) from e

        pix = page.get_pixmap(dpi=200)
        img_data = pix.tobytes("png")
        with Image.open(io.BytesIO(img_data)) as img:
            return _ocr_image(img)

    def extract(self, path: str) -> str:
        # Prefer PyMuPDF (fitz) for speed and native PDF text extraction.
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(path)
            parts: List[str] = []
            for page in doc:
                page_text = page.get_text("text") or ""
                if not page_text.strip() or _is_low_quality_text(page_text):
                    ocr_text = self._ocr_page(page)
                    if _text_quality_score(ocr_text) > _text_quality_score(page_text):
                        parts.append(ocr_text)
                    else:
                        parts.append(page_text)
                else:
                    parts.append(page_text)
            return "\n".join(parts)
        except Exception:
            # Fallback to pdfminer.six extraction only.
            try:
                from pdfminer.high_level import extract_text
            except Exception as e:
                raise RuntimeError(
                    "Missing dependency: pymupdf or pdfminer.six (install via requirements)"
                ) from e

            try:
                return extract_text(path) or ""
            except Exception:
                return ""


def file_metadata(path: str) -> dict:
    st = os.stat(path)
    return {
        "size": st.st_size,
        "mtime": int(st.st_mtime),
        "ctime": int(st.st_ctime),
        "basename": os.path.basename(path),
        "ext": os.path.splitext(path)[1].lower(),
    }
