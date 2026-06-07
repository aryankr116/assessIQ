from docx import Document
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
from pathlib import Path

base = Path(__file__).resolve().parent.parent / 'backend' / 'sample_docs'
base.mkdir(parents=True, exist_ok=True)

# DOCX
doc = Document()
doc.add_heading('Sample DOCX', level=1)
doc.add_paragraph('This document is used to test .docx ingestion.')
doc.add_paragraph('Highlights:')
doc.add_paragraph('- Word extraction\n- Paragraph handling')
doc.save(base / 'sample3.docx')

# PPTX
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Sample PPTX'
body = slide.shapes.add_textbox(left=1000000, top=1500000, width=8000000, height=2000000)
frame = body.text_frame
frame.text = 'This slide tests PowerPoint text extraction.'
prs.save(base / 'sample4.pptx')

# PNG with text
img = Image.new('RGB', (800, 400), color='white')
d = ImageDraw.Draw(img)
font = ImageFont.load_default()
d.text((20, 20), 'Sample Image OCR\nThis text should be detected if Tesseract is installed.', fill='black', font=font)
img.save(base / 'sample5.png')

print('created sample3.docx sample4.pptx sample5.png')
